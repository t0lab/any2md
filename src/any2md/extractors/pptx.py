"""pptx extractor — python-pptx based, pure-Python."""
from __future__ import annotations

import logging
import re
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


_EMU_PER_PX = 9525                                    # 96 DPI


def _emu_to_px(v: Any) -> float:
    if v is None:
        return 0.0
    try:
        return round(int(v) / _EMU_PER_PX, 1)
    except Exception:
        return 0.0

from any2md.extractors._smartart import (
    detect_layout_kind,
    parse_rels,
    parse_smartart_tree,
)
from any2md.extractors._image_utils import image_dimensions, normalize_to_png
from any2md.extractors.base import ExtractResult

log = logging.getLogger(__name__)

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}


class PptxExtractor:
    format = "pptx"

    def extract(self, source: Path) -> ExtractResult:
        warnings: list[str] = []
        images: dict[str, bytes] = {}
        prs = Presentation(source)
        slides_ir: list[dict[str, Any]] = []

        # Pre-load diagram XMLs + per-slide rels from the zip
        self._dgm_data_xmls: dict[str, bytes] = {}
        self._dgm_layout_xmls: dict[str, bytes] = {}
        # slide_rels[slide_idx] = {rel_id: target_path_relative_to_pptx_root}
        self._slide_rels: dict[int, dict[str, str]] = {}
        # comments per slide
        self._slide_comments: dict[int, list[dict[str, Any]]] = {}
        self._comment_authors: dict[str, str] = {}      # id → name

        try:
            with zipfile.ZipFile(source) as z:
                names = z.namelist()
                # First pass: load comment authors
                if "ppt/commentAuthors.xml" in names:
                    try:
                        self._comment_authors = self._parse_comment_authors(
                            z.read("ppt/commentAuthors.xml")
                        )
                    except Exception as e:
                        warnings.append(f"commentAuthors parse failed: {e}")

                # Cache comment XMLs by path for later resolution via slide rels
                comment_xmls: dict[str, bytes] = {}
                for name in names:
                    if name.startswith("ppt/comments/comment") and name.endswith(".xml"):
                        comment_xmls[name] = z.read(name)

                for name in names:
                    if name.startswith("ppt/diagrams/data") and name.endswith(".xml"):
                        self._dgm_data_xmls[name] = z.read(name)
                    elif name.startswith("ppt/diagrams/layout") and name.endswith(".xml"):
                        self._dgm_layout_xmls[name] = z.read(name)
                    else:
                        m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels$", name)
                        if m:
                            slide_idx = int(m.group(1))
                            try:
                                raw_rels = parse_rels(z.read(name))
                                rels_map: dict[str, str] = {}
                                for rid, target in raw_rels.items():
                                    if target.startswith("../"):
                                        target = "ppt/" + target[3:]
                                    elif target.startswith("/"):
                                        target = target.lstrip("/")
                                    rels_map[rid] = target
                                self._slide_rels[slide_idx] = rels_map
                                # Resolve any comments rel into this slide's comment list
                                for target in rels_map.values():
                                    if target.startswith("ppt/comments/") \
                                            and target.endswith(".xml") \
                                            and target in comment_xmls:
                                        try:
                                            cmts = self._parse_comments(
                                                comment_xmls[target]
                                            )
                                            self._slide_comments.setdefault(
                                                slide_idx, []
                                            ).extend(cmts)
                                        except Exception as e:
                                            warnings.append(
                                                f"Slide {slide_idx} comments parse failed: {e}"
                                            )
                            except Exception as e:
                                warnings.append(f"Slide {slide_idx} rels parse failed: {e}")
        except Exception as e:
            warnings.append(f"Could not pre-load zip contents: {e}")

        img_counter = 0
        for slide_idx, slide in enumerate(prs.slides, start=1):
            title = self._detect_title(slide)
            raw_shapes: list[dict[str, Any]] = []
            for shape in slide.shapes:
                ir_shape, img_counter, new_images, new_warnings = self._extract_shape(
                    shape, slide_idx, title, img_counter
                )
                if ir_shape is not None:
                    raw_shapes.append(ir_shape)
                images.update(new_images)
                warnings.extend(new_warnings)

            # Sort by (top, left) for reading order
            sorted_shapes = sorted(
                raw_shapes,
                key=lambda s: (s.get("bbox", {}).get("top", 0), s.get("bbox", {}).get("left", 0)),
            )
            for order, shp in enumerate(sorted_shapes):
                shp["order"] = order

            # Wire image neighbors based on order
            self._wire_image_neighbors(sorted_shapes)

            notes = None
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
            except Exception:
                pass

            slide_ir: dict[str, Any] = {
                "index": slide_idx,
                "title": title,
                "shapes": sorted_shapes,
                "notes": notes,
            }
            cmts = self._slide_comments.get(slide_idx)
            if cmts:
                slide_ir["comments"] = cmts
            slides_ir.append(slide_ir)

        ir = {
            "format": "pptx",
            "source": str(source),
            "slide_count": len(slides_ir),
            "slides": slides_ir,
        }
        return ExtractResult(ir=ir, images=images, warnings=warnings)

    # ----------------- shape extraction -----------------

    def _extract_shape(
        self,
        shape: Any,
        slide_idx: int,
        slide_label: str | None,
        img_counter: int,
    ) -> tuple[dict[str, Any] | None, int, dict[str, bytes], list[str]]:
        images: dict[str, bytes] = {}
        warnings: list[str] = []
        bbox = self._bbox(shape)

        try:
            shape_type = shape.shape_type
        except Exception:
            shape_type = None

        # SmartArt: detect via graphicFrame containing dgm: data
        if self._is_smartart(shape):
            ir, sa_warnings = self._extract_smartart(shape, bbox, slide_idx)
            warnings.extend(sa_warnings)
            return ir, img_counter, images, warnings

        # Chart
        try:
            if shape.has_chart:
                return self._extract_chart_shape(shape, bbox), img_counter, images, warnings
        except (AttributeError, ValueError):
            pass

        # Table
        try:
            if shape.has_table:
                return self._extract_table_shape(shape, bbox), img_counter, images, warnings
        except (AttributeError, ValueError):
            pass

        # Picture
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_counter += 1
            return self._extract_picture(
                shape, bbox, slide_idx, slide_label, img_counter, images, warnings
            )

        # Group
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            children: list[dict[str, Any]] = []
            for child in shape.shapes:
                c_ir, img_counter, c_images, c_warnings = self._extract_shape(
                    child, slide_idx, slide_label, img_counter
                )
                if c_ir is not None:
                    children.append(c_ir)
                images.update(c_images)
                warnings.extend(c_warnings)
            children.sort(
                key=lambda s: (s.get("bbox", {}).get("top", 0), s.get("bbox", {}).get("left", 0))
            )
            for order, c in enumerate(children):
                c["order"] = order
            diagram_hint = self._detect_flowchart(children)
            return (
                {
                    "type": "group",
                    "bbox": bbox,
                    "children": children,
                    "diagram_hint": diagram_hint,
                },
                img_counter,
                images,
                warnings,
            )

        # Connector
        if self._is_connector(shape):
            return self._extract_connector(shape, bbox), img_counter, images, warnings

        # Text
        try:
            if shape.has_text_frame:
                return self._extract_text_shape(shape, bbox), img_counter, images, warnings
        except (AttributeError, ValueError):
            pass

        # Unknown shape — emit minimal info + warning
        warnings.append(
            f"Slide {slide_idx}: unknown shape type {shape_type!r} (name={getattr(shape, 'name', '?')})"
        )
        return None, img_counter, images, warnings

    def _bbox(self, shape: Any) -> dict[str, float]:
        try:
            return {
                "left": _emu_to_px(shape.left),
                "top": _emu_to_px(shape.top),
                "width": _emu_to_px(shape.width),
                "height": _emu_to_px(shape.height),
                "unit": "px",
            }
        except Exception:
            return {"left": 0.0, "top": 0.0, "width": 0.0, "height": 0.0, "unit": "px"}

    def _detect_title(self, slide: Any) -> str | None:
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                return slide.shapes.title.text_frame.text.strip() or None
        except Exception:
            return None
        return None

    def _extract_text_shape(self, shape: Any, bbox: dict[str, float]) -> dict[str, Any]:
        paragraphs = [self._paragraph_to_ir(p) for p in shape.text_frame.paragraphs]
        return {
            "type": "text",
            "bbox": bbox,
            "paragraphs": paragraphs,
        }

    def _paragraph_to_ir(self, p: Any) -> dict[str, Any]:
        runs: list[dict[str, Any]] = []
        for r in p.runs:
            run_ir: dict[str, Any] = {"text": r.text or ""}
            try:
                if r.font.bold:
                    run_ir["bold"] = True
                if r.font.italic:
                    run_ir["italic"] = True
                if r.font.underline:
                    run_ir["underline"] = True
                if r.hyperlink and r.hyperlink.address:
                    run_ir["hyperlink"] = r.hyperlink.address
            except Exception:
                pass
            runs.append(run_ir)
        if not runs and p.text:
            runs.append({"text": p.text})
        para_ir: dict[str, Any] = {"runs": runs}
        try:
            if p.level:
                para_ir["level"] = int(p.level)
        except Exception:
            pass
        # Bullet kind + alignment from underlying <a:pPr>
        try:
            pPr = p._p.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
            )
            if pPr is not None:
                # buChar → bullet; buAutoNum → numbered; buNone → no bullet (skip)
                if pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar") is not None:
                    para_ir["list_kind"] = "bullet"
                elif pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum") is not None:
                    para_ir["list_kind"] = "numbered"
                algn = pPr.get("algn")
                if algn:
                    algn_map = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}
                    if algn in algn_map:
                        para_ir["alignment"] = algn_map[algn]
        except Exception:
            pass
        return para_ir

    def _extract_table_shape(self, shape: Any, bbox: dict[str, float]) -> dict[str, Any]:
        rows: list[list[str]] = []
        rows_rich: list[list[dict[str, Any]]] = []
        for row in shape.table.rows:
            flat_row: list[str] = []
            rich_row: list[dict[str, Any]] = []
            for cell in row.cells:
                flat_row.append(cell.text or "")
                cell_paragraphs: list[dict[str, Any]] = []
                try:
                    for p in cell.text_frame.paragraphs:
                        cell_paragraphs.append(self._paragraph_to_ir(p))
                except Exception:
                    pass
                rich_cell: dict[str, Any] = {
                    "paragraphs": cell_paragraphs,
                    "text": cell.text or "",
                }
                # Capture merge info if present
                try:
                    if getattr(cell, "is_merge_origin", False):
                        if cell.span_height > 1:
                            rich_cell["rowspan"] = int(cell.span_height)
                        if cell.span_width > 1:
                            rich_cell["colspan"] = int(cell.span_width)
                except Exception:
                    pass
                rich_row.append(rich_cell)
            rows.append(flat_row)
            rows_rich.append(rich_row)
        return {
            "type": "table",
            "bbox": bbox,
            "rows": rows,
            "rows_rich": rows_rich,
        }

    def _parse_comment_authors(self, raw: bytes) -> dict[str, str]:
        from lxml import etree
        try:
            tree = etree.fromstring(raw)
        except Exception:
            return {}
        ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        out: dict[str, str] = {}
        for ca in tree.findall(f".//{ns}cmAuthor"):
            aid = ca.get("id")
            name = ca.get("name") or ca.get("initials") or "Unknown"
            if aid is not None:
                out[aid] = name
        return out

    def _parse_comments(self, raw: bytes) -> list[dict[str, Any]]:
        from lxml import etree
        try:
            tree = etree.fromstring(raw)
        except Exception:
            return []
        ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        out: list[dict[str, Any]] = []
        for cm in tree.findall(f".//{ns}cm"):
            author_id = cm.get("authorId")
            author = self._comment_authors.get(author_id, None) if author_id else None
            text_el = cm.find(f"{ns}text")
            text = (text_el.text or "").strip() if text_el is not None else ""
            if not text:
                continue
            entry: dict[str, Any] = {"author": author, "text": text}
            dt = cm.get("dt")
            if dt:
                entry["created"] = dt
            out.append(entry)
        return out

    def _extract_picture(
        self,
        shape: Any,
        bbox: dict[str, float],
        slide_idx: int,
        slide_label: str | None,
        img_counter: int,
        images: dict[str, bytes],
        warnings: list[str],
    ) -> tuple[dict[str, Any], int, dict[str, bytes], list[str]]:
        try:
            blob = shape.image.blob
        except Exception as e:
            warnings.append(f"Slide {slide_idx}: picture extract failed: {e}")
            return None, img_counter, images, warnings  # type: ignore[return-value]

        image_id = f"img_s{slide_idx}_{img_counter}"
        rel_path = f"raw_images/{image_id}.png"
        normalized = normalize_to_png(blob)
        images[rel_path] = normalized
        w, h = image_dimensions(normalized)
        alt = ""
        try:
            alt = shape.element.nvSpPr.cNvPr.descr or ""
        except Exception:
            try:
                alt = shape._element.xpath("./p:nvPicPr/p:cNvPr/@descr")[0]
            except Exception:
                pass

        ir = {
            "type": "image",
            "bbox": bbox,
            "image": {
                "id": image_id,
                "relative_path": rel_path,
                "width_px": w,
                "height_px": h,
                "alt_text": alt or None,
                "bbox": bbox,
                "container": {
                    "kind": "pptx_slide",
                    "index": slide_idx,
                    "label": slide_label,
                },
                "neighbors": {
                    "before_text": "",
                    "after_text": "",
                    "before_kind": None,
                    "after_kind": None,
                },
            },
        }
        return ir, img_counter, images, warnings

    def _extract_chart_shape(self, shape: Any, bbox: dict[str, float]) -> dict[str, Any]:
        chart = shape.chart
        chart_type = "other"
        try:
            chart_type_name = str(chart.chart_type).lower()
            for key in ("bar", "column", "line", "pie", "scatter", "area", "doughnut"):
                if key in chart_type_name:
                    chart_type = key
                    break
        except Exception:
            pass

        title = None
        try:
            if chart.has_title and chart.chart_title.text_frame:
                title = chart.chart_title.text_frame.text
        except Exception:
            pass

        categories: list[str] = []
        series_ir: list[dict[str, Any]] = []
        try:
            for plot in chart.plots:
                if not categories:
                    try:
                        categories = [str(c) for c in plot.categories]
                    except Exception:
                        pass
                for ser in plot.series:
                    name = None
                    try:
                        name = ser.name
                    except Exception:
                        pass
                    values: list[Any] = []
                    try:
                        values = list(ser.values)
                    except Exception:
                        pass
                    series_ir.append({"name": name, "values": values})
        except Exception:
            pass

        return {
            "type": "chart",
            "bbox": bbox,
            "chart_type": chart_type,
            "title": title,
            "axis_labels": {},
            "categories": categories,
            "series": series_ir,
        }

    # ===================== SmartArt =====================

    def _is_smartart(self, shape: Any) -> bool:
        try:
            # SmartArt = graphicFrame with graphicData uri containing "diagram"
            elem = shape._element
            if elem.tag.split("}")[-1] != "graphicFrame":
                return False
            gd = elem.find(".//a:graphicData", namespaces=_NS)
            if gd is None:
                return False
            uri = gd.get("uri", "")
            return "diagram" in uri
        except Exception:
            return False

    def _extract_smartart(
        self, shape: Any, bbox: dict[str, float], slide_idx: int
    ) -> tuple[dict[str, Any], list[str]]:
        warnings: list[str] = []
        nodes: list[dict[str, Any]] = []
        layout_kind = "other"

        # 1) Resolve relIds → specific data and layout XML files
        data_xml_bytes: bytes | None = None
        layout_xml_bytes: bytes | None = None
        try:
            rel = shape._element.find(".//dgm:relIds", namespaces=_NS)
            if rel is not None and slide_idx in self._slide_rels:
                rels_map = self._slide_rels[slide_idx]
                dm_rid = rel.get(f"{{{_NS['r']}}}dm")
                lo_rid = rel.get(f"{{{_NS['r']}}}lo")
                if dm_rid and dm_rid in rels_map:
                    data_xml_bytes = self._dgm_data_xmls.get(rels_map[dm_rid])
                    if data_xml_bytes is None:
                        warnings.append(
                            f"Slide {slide_idx}: SmartArt data XML not found at "
                            f"{rels_map[dm_rid]!r}"
                        )
                if lo_rid and lo_rid in rels_map:
                    layout_xml_bytes = self._dgm_layout_xmls.get(rels_map[lo_rid])
                    if layout_xml_bytes is None:
                        warnings.append(
                            f"Slide {slide_idx}: SmartArt layout XML not found at "
                            f"{rels_map[lo_rid]!r}"
                        )
            else:
                warnings.append(
                    f"Slide {slide_idx}: SmartArt relIds element or slide rels missing — "
                    "using fallback (first diagram data file)"
                )
        except Exception as e:
            warnings.append(f"Slide {slide_idx}: SmartArt rels resolution failed: {e}")

        # Fallback if not resolved
        if data_xml_bytes is None and self._dgm_data_xmls:
            data_xml_bytes = next(iter(self._dgm_data_xmls.values()))
        if layout_xml_bytes is None and self._dgm_layout_xmls:
            layout_xml_bytes = next(iter(self._dgm_layout_xmls.values()))

        # 2) Parse the data XML → build node tree
        if data_xml_bytes is not None:
            try:
                nodes = parse_smartart_tree(data_xml_bytes)
            except Exception as e:
                warnings.append(f"Slide {slide_idx}: SmartArt data parse failed: {e}")

        # 3) Determine layout_kind from layout XML
        if layout_xml_bytes is not None:
            layout_kind = detect_layout_kind(layout_xml_bytes)

        return (
            {
                "type": "smartart",
                "bbox": bbox,
                "layout_kind": layout_kind,
                "nodes": nodes,
            },
            warnings,
        )

    # ===================== Connector =====================

    def _is_connector(self, shape: Any) -> bool:
        try:
            shape_type = shape.shape_type
        except Exception:
            return False
        if shape_type == MSO_SHAPE_TYPE.LINE:
            return True
        # Connector subtypes via xml tag
        try:
            tag = shape._element.tag.split("}")[-1]
            return tag == "cxnSp"
        except Exception:
            return False

    def _extract_connector(self, shape: Any, bbox: dict[str, float]) -> dict[str, Any]:
        from_id: str | None = None
        to_id: str | None = None
        try:
            stCxn = shape._element.xpath(".//a:stCxn", namespaces=_NS)
            endCxn = shape._element.xpath(".//a:endCxn", namespaces=_NS)
            if stCxn:
                from_id = stCxn[0].get("id")
            if endCxn:
                to_id = endCxn[0].get("id")
        except Exception:
            pass
        arrow_kind = "straight"
        try:
            tag = shape._element.tag.split("}")[-1]
            if tag == "cxnSp":
                prstGeom = shape._element.xpath(".//a:prstGeom", namespaces=_NS)
                if prstGeom:
                    preset = prstGeom[0].get("prst", "")
                    if "bent" in preset.lower():
                        arrow_kind = "elbow"
                    elif "curve" in preset.lower():
                        arrow_kind = "curved"
        except Exception:
            pass

        label = None
        try:
            if shape.has_text_frame:
                label = shape.text_frame.text.strip() or None
        except Exception:
            pass

        return {
            "type": "connector",
            "bbox": bbox,
            "from_shape_id": from_id,
            "to_shape_id": to_id,
            "arrow_kind": arrow_kind,
            "label": label,
        }

    def _detect_flowchart(self, children: list[dict[str, Any]]) -> str:
        connectors = [c for c in children if c.get("type") == "connector"]
        if len(connectors) >= 2 and all(
            c.get("from_shape_id") and c.get("to_shape_id") for c in connectors
        ):
            return "flowchart"
        return "none"

    # ===================== Image neighbors =====================

    def _wire_image_neighbors(self, shapes: list[dict[str, Any]]) -> None:
        for i, shp in enumerate(shapes):
            if shp.get("type") != "image":
                continue
            before_text = ""
            before_kind: str | None = None
            after_text = ""
            after_kind: str | None = None
            if i > 0:
                prev = shapes[i - 1]
                before_kind = prev.get("type")
                before_text = self._shape_to_text_snippet(prev)
            if i + 1 < len(shapes):
                nxt = shapes[i + 1]
                after_kind = nxt.get("type")
                after_text = self._shape_to_text_snippet(nxt)
            shp["image"]["neighbors"] = {
                "before_text": before_text[:200],
                "after_text": after_text[:200],
                "before_kind": before_kind,
                "after_kind": after_kind,
            }

    def _shape_to_text_snippet(self, shape: dict[str, Any]) -> str:
        t = shape.get("type")
        if t == "text":
            return " ".join(
                "".join(r.get("text", "") for r in p.get("runs", []))
                for p in shape.get("paragraphs", [])
            )
        if t == "table":
            return " | ".join(" ".join(row) for row in (shape.get("rows") or [])[:2])
        if t == "chart":
            return f"chart: {shape.get('title') or ''}"
        if t == "smartart":
            # Use first few node texts as snippet
            ns = shape.get("nodes", [])
            return "smartart: " + ", ".join(self._all_node_texts(ns)[:5])
        if t == "image":
            return "image"
        if t == "connector":
            return f"connector: {shape.get('label') or ''}"
        if t == "group":
            return "group"
        return ""

    def _all_node_texts(self, nodes: list[dict[str, Any]]) -> list[str]:
        out: list[str] = []
        for n in nodes:
            if n.get("text"):
                out.append(n["text"])
            out.extend(self._all_node_texts(n.get("children", [])))
        return out

